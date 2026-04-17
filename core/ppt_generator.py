import os

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
UPLOAD_DIR = os.getenv("UPLOAD_DIR", os.path.join(BASE_DIR, "uploads"))

PRESETS = {
    "corporate": {
        "name": "코퍼레이트",
        "desc": "격식 있는 기업 강의용. 안정적이고 신뢰감 있는 구성.",
        "colors": {
            "header": (0x1E, 0x3A, 0x5F),
            "accent": (0x4A, 0x90, 0xD9),
            "confirm": (0x1D, 0x9E, 0x75),
            "light": (0xEB, 0xF2, 0xFA),
            "bg": (0xFF, 0xFF, 0xFF),
            "text": (0x1A, 0x22, 0x33),
            "title_sub": (0xA8, 0xC8, 0xE8),
            "title_dark": (0x15, 0x2B, 0x4A),
        },
        "header_style": "full",
        "bullet_style": "circle",
        "density": "standard",
    },
    "startup": {
        "name": "스타트업",
        "desc": "대담하고 현대적인 스타트업 피치 스타일.",
        "colors": {
            "header": (0x0F, 0x17, 0x2A),
            "accent": (0x6C, 0x63, 0xFF),
            "confirm": (0x2D, 0xD4, 0xBF),
            "light": (0xF0, 0xEF, 0xFF),
            "bg": (0xFA, 0xFA, 0xFD),
            "text": (0x11, 0x18, 0x27),
            "title_sub": (0xA0, 0x9C, 0xFF),
            "title_dark": (0x07, 0x0D, 0x1C),
        },
        "header_style": "bottom_line",
        "bullet_style": "square",
        "density": "spacious",
    },
    "academic": {
        "name": "아카데믹",
        "desc": "학술/교육 자료. 정보 밀도 높고 깔끔한 레이아웃.",
        "colors": {
            "header": (0x2C, 0x3E, 0x50),
            "accent": (0xE7, 0x4C, 0x3C),
            "confirm": (0x27, 0xAE, 0x60),
            "light": (0xEC, 0xF0, 0xF1),
            "bg": (0xFF, 0xFF, 0xFF),
            "text": (0x1A, 0x1A, 0x1A),
            "title_sub": (0xBD, 0xC3, 0xC7),
            "title_dark": (0x1A, 0x28, 0x35),
        },
        "header_style": "left_bar",
        "bullet_style": "dash",
        "density": "compact",
    },
    "creative": {
        "name": "크리에이티브",
        "desc": "디자인/마케팅 분야. 생동감 있고 개성 강한 스타일.",
        "colors": {
            "header": (0xFF, 0x6B, 0x6B),
            "accent": (0xFF, 0xE6, 0x6D),
            "confirm": (0x4E, 0xCD, 0xC4),
            "light": (0xFF, 0xF5, 0xF5),
            "bg": (0xFF, 0xFF, 0xFF),
            "text": (0x2D, 0x2D, 0x2D),
            "title_sub": (0xFF, 0xC8, 0xC8),
            "title_dark": (0xCC, 0x44, 0x44),
        },
        "header_style": "full",
        "bullet_style": "number",
        "density": "spacious",
    },
    "terra": {
        "name": "테라코타",
        "desc": "따뜻하고 자연스러운 분위기. 인문/예술 강의에 적합.",
        "colors": {
            "header": (0x7A, 0x3B, 0x2E),
            "accent": (0xC8, 0x6B, 0x4A),
            "confirm": (0xD4, 0x9A, 0x3A),
            "light": (0xFD, 0xF0, 0xE8),
            "bg": (0xFF, 0xFB, 0xF7),
            "text": (0x2C, 0x1A, 0x12),
            "title_sub": (0xF0, 0xC4, 0xA8),
            "title_dark": (0x5A, 0x2A, 0x1E),
        },
        "header_style": "full",
        "bullet_style": "circle",
        "density": "standard",
    },
    "mono": {
        "name": "모노 미니멀",
        "desc": "흑백 모노톤. 텍스트 중심의 미니멀 디자인.",
        "colors": {
            "header": (0x18, 0x18, 0x18),
            "accent": (0x58, 0x58, 0x58),
            "confirm": (0x38, 0x38, 0x38),
            "light": (0xF2, 0xF2, 0xF2),
            "bg": (0xFF, 0xFF, 0xFF),
            "text": (0x1A, 0x1A, 0x1A),
            "title_sub": (0xB8, 0xB8, 0xB8),
            "title_dark": (0x0A, 0x0A, 0x0A),
        },
        "header_style": "bottom_line",
        "bullet_style": "dash",
        "density": "standard",
    },
    "forest": {
        "name": "포레스트",
        "desc": "자연/환경/웰니스. 차분하고 신뢰감 있는 초록 계열.",
        "colors": {
            "header": (0x1C, 0x40, 0x2E),
            "accent": (0x2E, 0x7D, 0x52),
            "confirm": (0x5A, 0xA8, 0x6A),
            "light": (0xE8, 0xF5, 0xED),
            "bg": (0xF7, 0xFC, 0xF8),
            "text": (0x12, 0x2A, 0x1C),
            "title_sub": (0xA0, 0xCC, 0xB4),
            "title_dark": (0x10, 0x2A, 0x1C),
        },
        "header_style": "full",
        "bullet_style": "circle",
        "density": "standard",
    },
    "pastel": {
        "name": "소프트 파스텔",
        "desc": "부드럽고 친근한 분위기. 아동/청소년 교육이나 워크숍.",
        "colors": {
            "header": (0x7E, 0x57, 0xC2),
            "accent": (0xF0, 0x6A, 0xA1),
            "confirm": (0x26, 0xC6, 0xDA),
            "light": (0xF3, 0xE5, 0xF5),
            "bg": (0xFF, 0xFF, 0xFF),
            "text": (0x31, 0x27, 0x3F),
            "title_sub": (0xD1, 0xC4, 0xE9),
            "title_dark": (0x51, 0x2D, 0xA8),
        },
        "header_style": "left_bar",
        "bullet_style": "circle",
        "density": "spacious",
    },
}

LEGACY_THEME_PRESET_MAP = {
    "navy": "corporate",
    "terra": "terra",
    "mono": "mono",
    "forest": "forest",
}

SUPPORTED_FONTS = [
    {"id": "Malgun Gothic", "label": "맑은 고딕"},
    {"id": "NanumGothic", "label": "나눔고딕"},
    {"id": "NanumBarunGothic", "label": "나눔바른고딕"},
    {"id": "NanumMyeongjo", "label": "나눔명조"},
    {"id": "Pretendard", "label": "프리텐다드"},
    {"id": "Apple SD Gothic Neo", "label": "애플 SD 고딕"},
]

THEMES = {
    "navy": {
        "name": "네이비 클래식",
        **{key: RGBColor(*value) for key, value in PRESETS["corporate"]["colors"].items()},
    },
    "terra": {
        "name": "따뜻한 테라코타",
        **{key: RGBColor(*value) for key, value in PRESETS["terra"]["colors"].items()},
    },
    "mono": {
        "name": "모노 미니멀",
        **{key: RGBColor(*value) for key, value in PRESETS["mono"]["colors"].items()},
    },
    "forest": {
        "name": "포레스트 그린",
        **{key: RGBColor(*value) for key, value in PRESETS["forest"]["colors"].items()},
    },
}

DENSITY_CONFIG = {
    "compact": {"gap": 0.92, "size": 16},
    "standard": {"gap": 1.04, "size": 18},
    "spacious": {"gap": 1.12, "size": 19},
}

CARD_COLS = 3
CARD_W = 4.0
CARD_H = 1.6
CARD_GAP = 0.2
CARD_START_X = 0.37
VALID_CONTENT_LAYOUTS = {"auto", "classic", "split", "card", "highlight", "process", "compare", "image_left", "image_top", "chapter"}
DEFAULT_TITLE_FONT_SIZE = 32
DEFAULT_SUBTITLE_FONT_SIZE = 18
DEFAULT_BODY_FONT_SIZE = 18


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _rect(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    _solid(shape, color)
    return shape


def _oval(slide, x, y, w, h, color):
    shape = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(w), Inches(h))
    _solid(shape, color)
    return shape


def _tb(
    slide,
    x,
    y,
    w,
    h,
    text,
    size,
    bold=False,
    color=None,
    align=PP_ALIGN.LEFT,
    wrap=True,
    font_name="Malgun Gothic",
):
    color = color or RGBColor(0x1A, 0x22, 0x33)
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.clear()
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = str(text or "")
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name or "Malgun Gothic"
    return txb


def _estimate_text_lines(text, width: float, size: int) -> int:
    raw = str(text or "").strip()
    if not raw:
        return 1

    per_line = max(int((width * 72) / max(size * 0.92, 1)), 8)
    total_lines = 0
    for line in raw.splitlines() or [raw]:
        cleaned = line.strip() or " "
        total_lines += max(1, (len(cleaned) + per_line - 1) // per_line)
    return min(max(total_lines, 1), 8)


def _estimate_text_height(text, width: float, size: int, minimum: float = 0.46) -> float:
    lines = _estimate_text_lines(text, width, size)
    line_height = max((size * 1.35) / 72, 0.22)
    return max(minimum, lines * line_height + 0.06)


def _responsive_text_size(text, base_size: int, minimum: int, step: int = 22) -> int:
    penalty = max(0, (len(str(text or "").strip()) - 40) // step)
    return max(minimum, base_size - penalty)


def _bullet_spacing(metrics: dict) -> float:
    return max(0.18, min(metrics["gap"] - 0.54, 0.42))


def _set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _hex_to_rgb(hex_str: str) -> RGBColor:
    value = str(hex_str or "").strip().lstrip("#")
    if len(value) != 6:
        raise ValueError("hex color must be 6 characters")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _coerce_font_size(value, default: int, minimum: int, maximum: int) -> int:
    try:
        size = int(str(value).strip())
    except (TypeError, ValueError, AttributeError):
        return default
    return max(minimum, min(maximum, size))


def _normalize_preset_id(preset_id: str) -> str:
    return LEGACY_THEME_PRESET_MAP.get(preset_id, preset_id or "corporate")


def _find_logo_path(logo_uid):
    if not logo_uid:
        return None
    for ext in ("png", "jpg", "jpeg", "gif", "webp"):
        candidate = os.path.join(UPLOAD_DIR, f"logo_{logo_uid}.{ext}")
        if os.path.exists(candidate):
            return os.path.abspath(candidate)
    return None


def _resolve_theme(design: dict) -> dict:
    design = design or {}
    preset_id = _normalize_preset_id(design.get("preset", "corporate"))
    preset = PRESETS.get(preset_id, PRESETS["corporate"])

    t = {key: RGBColor(*value) for key, value in preset["colors"].items()}

    primary_color = design.get("primary_color")
    if primary_color:
        try:
            color = _hex_to_rgb(primary_color)
            t["header"] = color
            t["title_dark"] = color
        except ValueError:
            pass

    accent_color = design.get("accent_color")
    if accent_color:
        try:
            color = _hex_to_rgb(accent_color)
            t["accent"] = color
            t["confirm"] = color
        except ValueError:
            pass

    content_layout = _normalize_layout_id(design.get("content_layout", "auto"))

    t["preset_id"] = preset_id
    t["preset_name"] = preset["name"]
    t["header_style"] = preset.get("header_style", "full")
    t["bullet_style"] = preset.get("bullet_style", "circle")
    t["density"] = preset.get("density", "standard")
    t["content_layout"] = content_layout
    t["font_name"] = design.get("font_name", "Malgun Gothic") or "Malgun Gothic"
    t["company_name"] = design.get("company_name", "") or ""
    t["presenter_name"] = design.get("presenter_name", "") or ""
    t["footer_text"] = design.get("footer_text", "") or ""
    t["footer_enabled"] = design.get("footer_enabled", True) is not False
    t["slide_number"] = design.get("slide_number", True) is not False
    t["title_font_size"] = _coerce_font_size(
        design.get("title_font_size"),
        DEFAULT_TITLE_FONT_SIZE,
        20,
        56,
    )
    t["subtitle_font_size"] = _coerce_font_size(
        design.get("subtitle_font_size"),
        DEFAULT_SUBTITLE_FONT_SIZE,
        12,
        36,
    )
    t["body_font_size"] = _coerce_font_size(
        design.get("body_font_size"),
        DEFAULT_BODY_FONT_SIZE,
        12,
        28,
    )
    t["cover_title_font_size"] = min(t["title_font_size"] + 10, 64)
    t["cover_subtitle_font_size"] = min(t["subtitle_font_size"] + 4, 40)
    t["logo_path"] = _find_logo_path(design.get("logo_uid"))
    return t


def _normalize_layout_id(layout_id) -> str:
    layout = str(layout_id or "").strip().lower()
    if layout in VALID_CONTENT_LAYOUTS:
        return layout
    return "auto"


def _infer_auto_layout(data: dict) -> str:
    if str(data.get("role") or "").lower() == "chapter":
        return "chapter"

    points = [str(point).strip() for point in data.get("points", []) if str(point or "").strip()]
    subtitle = str(data.get("subtitle", "") or "").strip()
    image_mode = str(data.get("image_mode") or "none").strip().lower()
    orientation = str(data.get("image_orientation") or "square").strip().lower()

    if not points:
        return "classic"

    if image_mode == "hero" and data.get("image_asset_name"):
        return "image_left" if orientation == "portrait" else "image_top"

    lengths = [len(point) for point in points]
    max_len = max(lengths)
    first_len = lengths[0]
    short_points = sum(1 for length in lengths if length <= 36)
    remaining_lengths = lengths[1:]
    remaining_avg = (sum(remaining_lengths) / len(remaining_lengths)) if remaining_lengths else 0

    if subtitle and len(points) <= 4:
        return "split"

    if len(points) >= 4 and short_points >= min(len(points), 4) and max_len <= 60:
        return "card"

    if len(points) <= 2:
        return "highlight"

    if first_len >= 42 and (len(points) <= 4 or first_len >= remaining_avg + 12):
        return "highlight"

    if len(points) == 3 and subtitle:
        return "split"

    return "classic"


def _resolve_content_layout(data: dict, t: dict) -> str:
    if str(data.get("role") or "").lower() == "chapter":
        return "chapter"

    slide_layout = _normalize_layout_id(data.get("layout"))
    theme_layout = _normalize_layout_id(t.get("content_layout"))

    if slide_layout != "auto":
        return slide_layout

    if theme_layout != "auto":
        return theme_layout

    return _infer_auto_layout(data)


def _body_metrics(t: dict) -> dict:
    density_key = t.get("density", "standard")
    density = DENSITY_CONFIG.get(density_key, DENSITY_CONFIG["standard"])
    body_size = t.get("body_font_size", density["size"])
    extra_gap = max(body_size - density["size"], 0) * 0.06
    return {
        "size": body_size,
        "gap": density["gap"] + extra_gap,
    }


def _render_header(slide, data, t):
    title = data.get("title", "")
    header_style = t.get("header_style", "full")
    font_name = t.get("font_name", "Malgun Gothic")
    title_size = t.get("title_font_size", DEFAULT_TITLE_FONT_SIZE)

    if header_style == "left_bar":
        _rect(slide, 0, 0, 0.22, 7.5, t["accent"])
        _tb(
            slide,
            0.55,
            0.32,
            11.9,
            0.72,
            title,
            title_size,
            bold=True,
            color=t["header"],
            font_name=font_name,
        )
        return {"body_left": 0.58, "body_top": 1.25, "body_width": 12.1}

    if header_style == "bottom_line":
        _tb(
            slide,
            0.5,
            0.32,
            12.1,
            0.72,
            title,
            title_size,
            bold=True,
            color=t["header"],
            font_name=font_name,
        )
        _rect(slide, 0.5, 1.1, 12.33, 0.05, t["accent"])
        return {"body_left": 0.5, "body_top": 1.35, "body_width": 12.0}

    _rect(slide, 0, 0, 13.33, 1.28, t["header"])
    _rect(slide, 0, 1.28, 13.33, 0.065, t["accent"])
    _tb(
        slide,
        0.5,
        0.18,
        12.3,
        1.0,
        title,
        title_size,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        font_name=font_name,
    )
    return {"body_left": 0.5, "body_top": 1.62, "body_width": 12.0}


def _add_bullet(slide, x, y, pt, t, index, width=12.0, size=15):
    style = t.get("bullet_style", "circle")
    font_name = t.get("font_name", "Malgun Gothic")
    bullet_x = x
    text_x = x + 0.38

    if style == "square":
        _rect(slide, bullet_x, y + 0.10, 0.18, 0.18, t["accent"])
        text_x = x + 0.34
    elif style == "number":
        _rect(slide, bullet_x, y + 0.04, 0.30, 0.30, t["accent"])
        _tb(
            slide,
            bullet_x,
            y + 0.03,
            0.30,
            0.32,
            str(index + 1),
            max(size - 4, 10),
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            font_name=font_name,
        )
        text_x = x + 0.42
    elif style == "dash":
        _rect(slide, bullet_x, y + 0.18, 0.20, 0.04, t["accent"])
        text_x = x + 0.34
    else:
        _oval(slide, bullet_x, y + 0.09, 0.21, 0.21, t["accent"])

    text_width = max(width - (text_x - x), 0.5)
    text_height = _estimate_text_height(pt, text_width, size, minimum=0.48)
    _tb(slide, text_x, y, text_width, text_height, pt, size, color=t["text"], font_name=font_name)
    return text_height


def _add_footer(slide, slide_index: int, t: dict):
    footer_size = max(t.get("body_font_size", DEFAULT_BODY_FONT_SIZE) - 8, 10)

    if t.get("footer_enabled", True):
        _rect(slide, 0, 6.88, 13.33, 0.62, t["light"])

        if t.get("footer_text"):
            _tb(
                slide,
                0.4,
                6.9,
                8.0,
                0.45,
                t["footer_text"],
                footer_size,
                color=t["accent"],
                align=PP_ALIGN.LEFT,
                font_name=t.get("font_name"),
            )

        if t.get("slide_number", True):
            _tb(
                slide,
                11.0,
                6.9,
                2.0,
                0.45,
                str(slide_index),
                footer_size,
                color=t["text"],
                align=PP_ALIGN.RIGHT,
                font_name=t.get("font_name"),
            )
        return

    if t.get("slide_number", True):
        _tb(
            slide,
            11.2,
            6.7,
            1.6,
            0.35,
            str(slide_index),
            footer_size,
            color=t["text"],
            align=PP_ALIGN.RIGHT,
            font_name=t.get("font_name"),
        )


def _add_logo(slide, logo_path: str):
    if not logo_path or not os.path.exists(logo_path):
        return
    try:
        with Image.open(logo_path) as img:
            img_w, img_h = img.size
        max_w = Inches(1.1)
        max_h = Inches(0.78)
        ratio = min(max_w / img_w, max_h / img_h)
        width = int(img_w * ratio)
        height = int(img_h * ratio)
        x = Inches(13.33) - width - Inches(0.15)
        y = Inches(0.12)
        slide.shapes.add_picture(logo_path, x, y, width, height)
    except Exception:
        pass


def _set_notes(slide, data):
    notes = data.get("notes", "")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _slide_image_key(data: dict):
    return str(data.get("image_bundle_uid") or ""), str(data.get("image_asset_name") or "")


def _resolve_slide_image_asset(data: dict, media_assets: list[dict] | None):
    bundle_uid, asset_name = _slide_image_key(data)
    if not bundle_uid or not asset_name or not media_assets:
        return None
    for asset in media_assets:
        if asset.get("bundle_uid") == bundle_uid and asset.get("asset_name") == asset_name:
            return asset
    return None


def _image_orientation(image_asset: dict | None) -> str:
    if not image_asset:
        return "none"
    width = float(image_asset.get("width") or 0)
    height = float(image_asset.get("height") or 0)
    if width <= 0 or height <= 0:
        return "square"
    ratio = width / height
    if ratio >= 1.3:
        return "landscape"
    if ratio <= 0.82:
        return "portrait"
    return "square"


def _image_is_prominent(image_asset: dict | None, minimum_coverage: float = 0.014) -> bool:
    if not image_asset:
        return False
    coverage = float(image_asset.get("coverage_ratio") or 0)
    display_area = float(image_asset.get("display_area") or 0)
    return coverage >= minimum_coverage or display_area >= 26000


def _image_mode(data: dict) -> str:
    mode = str(data.get("image_mode") or "none").strip().lower()
    return mode if mode in {"hero", "support", "none"} else "none"


def _add_content_image(slide, image_asset: dict | None, x: float, y: float, w: float, h: float):
    path = image_asset.get("path") if image_asset else None
    if not path or not os.path.exists(path):
        return False

    try:
        with Image.open(path) as img:
            img_w, img_h = img.size
        max_w = Inches(w)
        max_h = Inches(h)
        ratio = min(max_w / img_w, max_h / img_h)
        width = int(img_w * ratio)
        height = int(img_h * ratio)
        x_pos = Inches(x) + int((Inches(w) - width) / 2)
        y_pos = Inches(y) + int((Inches(h) - height) / 2)
        slide.shapes.add_picture(path, x_pos, y_pos, width, height)
        return True
    except Exception:
        return False


def _title_slide(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["header"])
    _rect(slide, 0, 5.9, 13.33, 1.6, t["title_dark"])
    _rect(slide, 0.8, 3.05, 5.0, 0.06, t["accent"])
    _tb(
        slide,
        0.8,
        1.2,
        11.7,
        1.6,
        data.get("title", "강의 교안"),
        t.get("cover_title_font_size", 42),
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        font_name=t["font_name"],
    )

    subtitle = data.get("subtitle", "")
    if subtitle:
        _tb(
            slide,
            0.8,
            3.2,
            11.7,
            0.9,
            subtitle,
            t.get("cover_subtitle_font_size", 22),
            color=t["title_sub"],
            font_name=t["font_name"],
        )

    if t.get("company_name"):
        _tb(
            slide,
            0.8,
            6.15,
            11.5,
            0.55,
            t["company_name"],
            max(t.get("body_font_size", DEFAULT_BODY_FONT_SIZE) - 2, 12),
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.RIGHT,
            font_name=t["font_name"],
        )

    if t.get("presenter_name"):
        _tb(
            slide,
            0.8,
            6.65,
            11.5,
            0.45,
            t["presenter_name"],
            max(t.get("subtitle_font_size", DEFAULT_SUBTITLE_FONT_SIZE) - 4, 11),
            color=t["title_sub"],
            align=PP_ALIGN.RIGHT,
            font_name=t["font_name"],
        )

    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _chapter_slide(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["header"])
    _rect(slide, 0, 5.95, 13.33, 1.55, t["title_dark"])
    _rect(slide, 0.72, 3.32, 2.2, 0.08, t["accent"])
    _tb(
        slide,
        0.72,
        1.35,
        8.4,
        1.8,
        data.get("title", "새 챕터"),
        min(t.get("cover_title_font_size", 42) + 2, 56),
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        font_name=t["font_name"],
    )
    subtitle = data.get("subtitle", "")
    if subtitle:
        _tb(
            slide,
            0.75,
            3.56,
            7.0,
            0.55,
            subtitle,
            max(t.get("subtitle_font_size", DEFAULT_SUBTITLE_FONT_SIZE), 16),
            color=t["title_sub"],
            font_name=t["font_name"],
        )

    points = [str(point).strip() for point in data.get("points", []) if str(point).strip()]
    if points:
        _rect(slide, 0.74, 4.52, 8.2, 0.92, t["light"])
        _tb(
            slide,
            0.98,
            4.75,
            7.68,
            0.42,
            points[0],
            max(t.get("body_font_size", DEFAULT_BODY_FONT_SIZE) - 1, 14),
            color=t["text"],
            font_name=t["font_name"],
        )

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _agenda_slide(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg"])
    layout = _render_header(slide, data, t)
    metrics = _body_metrics(t)
    y = layout["body_top"] + 0.05
    spacing = max(0.14, min(metrics["gap"] - 0.48, 0.34))
    bottom_limit = 6.55 if t.get("footer_enabled", True) else 6.88

    for i, item in enumerate(data.get("items", [])[:8]):
        if y >= bottom_limit:
            break
        _rect(slide, layout["body_left"], y + 0.04, 0.32, 0.32, t["accent"])
        _tb(
            slide,
            layout["body_left"],
            y + 0.03,
            0.32,
            0.33,
            str(i + 1),
            12,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            font_name=t["font_name"],
        )
        _tb(
            slide,
            layout["body_left"] + 0.48,
            y,
            11.5,
            _estimate_text_height(item, 11.5, t.get("body_font_size", DEFAULT_BODY_FONT_SIZE), minimum=0.48),
            item,
            t.get("body_font_size", DEFAULT_BODY_FONT_SIZE),
            color=t["text"],
            font_name=t["font_name"],
        )
        y += _estimate_text_height(item, 11.02, t.get("body_font_size", DEFAULT_BODY_FONT_SIZE), minimum=0.48) + spacing

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _content_slide_classic(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg"])
    layout = _render_header(slide, data, t)
    density = _body_metrics(t)
    image_asset = _resolve_slide_image_asset(data, media_assets)

    body_left = layout["body_left"]
    body_top = layout["body_top"]
    body_width = layout["body_width"]
    image_drawn = False
    if _image_mode(data) in {"support", "hero"} and _image_is_prominent(image_asset):
        orientation = _image_orientation(image_asset)
        if orientation == "portrait":
            image_drawn = _add_content_image(slide, image_asset, 8.54, max(body_top - 0.02, 1.5), 4.08, 4.92)
            if image_drawn:
                body_width = 7.02
        elif orientation == "square":
            image_drawn = _add_content_image(slide, image_asset, 8.2, max(body_top, 1.54), 4.35, 4.2)
            if image_drawn:
                body_width = 6.8
        else:
            image_drawn = _add_content_image(slide, image_asset, 7.95, max(body_top + 0.1, 1.62), 4.7, 3.4)
            if image_drawn:
                body_width = 6.62

    y = body_top
    spacing = _bullet_spacing(density)
    bottom_limit = 6.48 if t.get("footer_enabled", True) else 6.82

    for i, pt in enumerate(data.get("points", [])[:5]):
        bullet_height = _add_bullet(
            slide,
            body_left,
            y,
            pt,
            t,
            i,
            width=body_width,
            size=density["size"],
        )
        y += bullet_height + spacing
        if y >= bottom_limit:
            break

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _content_slide_split(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg"])
    image_asset = _resolve_slide_image_asset(data, media_assets)
    _rect(slide, 0, 0, 4.5, 7.5, t["header"])
    _rect(slide, 4.5, 0.5, 0.05, 6.15, t["accent"])
    _tb(
        slide,
        0.38,
        1.75,
        3.55,
        1.4,
        data.get("title", ""),
        min(t.get("title_font_size", DEFAULT_TITLE_FONT_SIZE) + 2, 42),
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        font_name=t["font_name"],
    )

    subtitle = data.get("subtitle", "")
    if subtitle:
        _tb(
            slide,
            0.38,
            3.35,
            3.55,
            0.7,
            subtitle,
            max(t.get("subtitle_font_size", DEFAULT_SUBTITLE_FONT_SIZE) - 1, 12),
            color=t["title_sub"],
            font_name=t["font_name"],
        )

    density = _body_metrics(t)
    bullet_start_y = 1.2
    bullet_width = 7.7
    if _image_mode(data) in {"support", "hero"} and _image_is_prominent(image_asset):
        orientation = _image_orientation(image_asset)
        if orientation == "portrait":
            if _add_content_image(slide, image_asset, 8.78, 0.96, 3.52, 4.8):
                bullet_width = 3.74
                bullet_start_y = 1.08
        elif _add_content_image(slide, image_asset, 4.86, 0.88, 7.82, 2.48):
            bullet_start_y = 3.7
            bullet_width = 7.6

    y = bullet_start_y
    spacing = _bullet_spacing(density)
    bottom_limit = 6.48 if t.get("footer_enabled", True) else 6.82

    for i, pt in enumerate(data.get("points", [])[:5]):
        bullet_height = _add_bullet(slide, 4.9, y, pt, t, i, width=bullet_width, size=density["size"])
        y += bullet_height + spacing
        if y >= bottom_limit:
            break

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _content_slide_card(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg"])
    layout = _render_header(slide, data, t)
    pts = data.get("points", [])[:5]
    top_y = max(layout["body_top"] + 0.02, 1.55)
    grid_x = 0.55
    grid_y = top_y
    grid_width = 12.23
    cols = 3
    gap = 0.22
    card_height = 1.56

    card_width = (grid_width - gap * (cols - 1)) / cols

    for i, pt in enumerate(pts):
        col = i % cols
        row = i // cols
        x = grid_x + col * (card_width + gap)
        y = grid_y + row * (card_height + gap)
        _rect(slide, x, y, card_width, card_height, t["light"])
        _rect(slide, x, y, card_width, 0.07, t["accent"])
        _tb(
            slide,
            x + 0.18,
            y + 0.22,
            card_width - 0.34,
            card_height - 0.38,
            pt,
            _responsive_text_size(
                pt,
                max(t.get("body_font_size", DEFAULT_BODY_FONT_SIZE) - 3, 13),
                11,
                step=16,
            ),
            color=t["text"],
            font_name=t["font_name"],
        )

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _content_slide_highlight(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg"])
    layout = _render_header(slide, data, t)
    pts = data.get("points", [])[:5]
    highlight_y = max(layout["body_top"] + 0.02, 1.55)

    if pts:
        highlight_w = 12.33
        text_x = 0.75
        text_w = 11.8
        _rect(slide, 0.5, highlight_y, highlight_w, 2.0, t["accent"])
        _tb(
            slide,
            text_x,
            highlight_y + 0.2,
            text_w,
            1.6,
            pts[0],
            min(t.get("title_font_size", DEFAULT_TITLE_FONT_SIZE) + 4, 48),
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            font_name=t["font_name"],
        )

    remaining = pts[1:4]
    col_w = 12.33 / max(len(remaining), 1)
    for i, pt in enumerate(remaining):
        x = 0.5 + i * col_w
        _rect(slide, x + 0.05, highlight_y + 2.35, col_w - 0.15, 1.4, t["light"])
        _tb(
            slide,
            x + 0.2,
            highlight_y + 2.45,
            col_w - 0.3,
            1.2,
            pt,
            _responsive_text_size(
                pt,
                max(t.get("body_font_size", DEFAULT_BODY_FONT_SIZE) - 2, 13),
                11,
                step=18,
            ),
            color=t["text"],
            font_name=t["font_name"],
        )

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _content_slide_process(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg"])
    layout = _render_header(slide, data, t)
    steps = data.get("diagram_steps") or data.get("points", [])
    steps = [str(step).strip() for step in steps if str(step).strip()][:4]
    top_y = max(layout["body_top"], 1.55)
    start_x = 0.55
    total_width = 12.1
    extra_right = 11.95

    usable_steps = steps or [point for point in data.get("points", [])[:3] if str(point).strip()]
    step_count = max(len(usable_steps), 1)
    gap = 0.18
    step_width = min((total_width - gap * (step_count - 1)) / step_count, 2.8)
    box_y = top_y + 0.35
    box_h = 1.5

    for idx, step in enumerate(usable_steps):
        x = start_x + idx * (step_width + gap)
        _rect(slide, x, box_y, step_width, box_h, t["light"])
        _rect(slide, x, box_y, step_width, 0.08, t["accent"])
        _tb(
            slide,
            x + 0.18,
            box_y + 0.18,
            0.42,
            0.32,
            str(idx + 1),
            max(t.get("body_font_size", DEFAULT_BODY_FONT_SIZE) - 4, 12),
            bold=True,
            color=t["accent"],
            font_name=t["font_name"],
        )
        _tb(
            slide,
            x + 0.18,
            box_y + 0.5,
            step_width - 0.3,
            _estimate_text_height(
                step,
                step_width - 0.3,
                _responsive_text_size(
                    step,
                    max(t.get("body_font_size", DEFAULT_BODY_FONT_SIZE) - 1, 14),
                    11,
                    step=14,
                ),
                minimum=0.64,
            ),
            step,
            _responsive_text_size(
                step,
                max(t.get("body_font_size", DEFAULT_BODY_FONT_SIZE) - 1, 14),
                11,
                step=14,
            ),
            bold=True,
            color=t["text"],
            font_name=t["font_name"],
        )
        if idx < step_count - 1:
            arrow_x = x + step_width + 0.03
            _rect(slide, arrow_x, box_y + 0.66, gap - 0.06, 0.06, t["accent"])
            _tb(
                slide,
                arrow_x + 0.02,
                box_y + 0.49,
                0.22,
                0.32,
                "→",
                max(t.get("body_font_size", DEFAULT_BODY_FONT_SIZE) - 2, 13),
                bold=True,
                color=t["accent"],
                font_name=t["font_name"],
            )

    metrics = _body_metrics(t)
    extra_points = data.get("points", [])[len(usable_steps):]
    y = box_y + box_h + 0.26
    spacing = _bullet_spacing(metrics)
    for idx, point in enumerate(extra_points[:3]):
        bullet_height = _add_bullet(slide, 0.7, y, point, t, idx, width=extra_right - 0.7, size=max(metrics["size"] - 1, 14))
        y += bullet_height + spacing

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _content_slide_compare(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg"])
    layout = _render_header(slide, data, t)
    top_y = max(layout["body_top"], 1.55)
    left_title = data.get("compare_left_title") or "핵심 A"
    right_title = data.get("compare_right_title") or "핵심 B"
    left_points = data.get("compare_left_points") or data.get("points", [])[:2]
    right_points = data.get("compare_right_points") or data.get("points", [])[2:4]
    left_box_x = 0.6
    right_box_x = 6.95
    box_width = 5.7
    divider_x = 6.58
    body_y = top_y + 1.15
    bottom_limit = 6.42 if t.get("footer_enabled", True) else 6.78

    _rect(slide, left_box_x, top_y, box_width, 0.9, t["light"])
    _rect(slide, right_box_x, top_y, box_width, 0.9, t["light"])
    _tb(
        slide,
        left_box_x + 0.22,
        top_y + 0.18,
        box_width - 0.45,
        0.45,
        left_title,
        max(t.get("subtitle_font_size", DEFAULT_SUBTITLE_FONT_SIZE) + 1, 16),
        bold=True,
        color=t["accent"],
        font_name=t["font_name"],
    )
    _tb(
        slide,
        right_box_x + 0.23,
        top_y + 0.18,
        box_width - 0.45,
        0.45,
        right_title,
        max(t.get("subtitle_font_size", DEFAULT_SUBTITLE_FONT_SIZE) + 1, 16),
        bold=True,
        color=t["accent"],
        font_name=t["font_name"],
    )
    _rect(slide, divider_x, top_y + 0.08, 0.05, max(bottom_limit - top_y - 0.24, 3.6), t["accent"])

    metrics = _body_metrics(t)
    left_y = body_y
    right_y = body_y
    spacing = _bullet_spacing(metrics)
    for idx, point in enumerate(left_points[:3]):
        bullet_height = _add_bullet(slide, left_box_x + 0.22, left_y, point, t, idx, width=box_width - 0.45, size=max(metrics["size"] - 1, 14))
        left_y += bullet_height + spacing

    for idx, point in enumerate(right_points[:3]):
        bullet_height = _add_bullet(slide, right_box_x + 0.23, right_y, point, t, idx, width=box_width - 0.45, size=max(metrics["size"] - 1, 14))
        right_y += bullet_height + spacing

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _content_slide_image_left(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg"])
    layout = _render_header(slide, data, t)
    image_asset = _resolve_slide_image_asset(data, media_assets)
    body_top = max(layout["body_top"], 1.58)
    image_left = 0.45
    image_width = 6.0
    text_left = 6.7
    text_width = 5.95
    image_height = 5.08 if _image_orientation(image_asset) == "portrait" else 4.32
    image_y = body_top - 0.02

    if not _add_content_image(slide, image_asset, image_left, image_y, image_width, image_height):
        body_left = layout["body_left"]
        body_top = layout["body_top"]
        body_width = layout["body_width"]
        metrics = _body_metrics(t)
        y = body_top
        spacing = _bullet_spacing(metrics)
        bottom_limit = 6.46 if t.get("footer_enabled", True) else 6.82
        for i, pt in enumerate(data.get("points", [])[:5]):
            bullet_height = _add_bullet(slide, body_left, y, pt, t, i, width=body_width, size=metrics["size"])
            y += bullet_height + spacing
            if y >= bottom_limit:
                break
        _add_footer(slide, slide_index, t)
        _add_logo(slide, t.get("logo_path"))
        _set_notes(slide, data)
        return

    subtitle = data.get("subtitle", "")
    if subtitle:
        _tb(
            slide,
            text_left,
            body_top,
            text_width,
            0.42,
            subtitle,
            max(t.get("subtitle_font_size", DEFAULT_SUBTITLE_FONT_SIZE) - 1, 14),
            color=t["accent"],
            font_name=t["font_name"],
        )

    metrics = _body_metrics(t)
    y = body_top + (0.52 if subtitle else 0.0)
    spacing = _bullet_spacing(metrics)
    bottom_limit = 6.46 if t.get("footer_enabled", True) else 6.82
    point_limit = 4 if _image_orientation(image_asset) == "portrait" else 3
    for i, pt in enumerate(data.get("points", [])[:point_limit]):
        bullet_height = _add_bullet(slide, text_left, y, pt, t, i, width=text_width, size=max(metrics["size"] - 1, 15))
        y += bullet_height + spacing
        if y >= bottom_limit:
            break

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _content_slide_image_top(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["bg"])
    layout = _render_header(slide, data, t)
    image_asset = _resolve_slide_image_asset(data, media_assets)
    body_top = max(layout["body_top"], 1.58)
    image_height = 3.18 if _image_orientation(image_asset) == "landscape" else 2.92

    if not _add_content_image(slide, image_asset, 0.56, body_top - 0.02, 12.18, image_height):
        body_left = layout["body_left"]
        body_width = layout["body_width"]
        metrics = _body_metrics(t)
        y = body_top
        spacing = _bullet_spacing(metrics)
        bottom_limit = 6.44 if t.get("footer_enabled", True) else 6.82
        for i, pt in enumerate(data.get("points", [])[:5]):
            bullet_height = _add_bullet(slide, body_left, y, pt, t, i, width=body_width, size=metrics["size"])
            y += bullet_height + spacing
            if y >= bottom_limit:
                break
        _add_footer(slide, slide_index, t)
        _add_logo(slide, t.get("logo_path"))
        _set_notes(slide, data)
        return

    subtitle = data.get("subtitle", "")
    text_top = body_top + image_height + 0.18
    if subtitle:
        _tb(
            slide,
            0.62,
            text_top,
            12.02,
            0.42,
            subtitle,
            max(t.get("subtitle_font_size", DEFAULT_SUBTITLE_FONT_SIZE) - 1, 14),
            color=t["accent"],
            font_name=t["font_name"],
        )
        text_top += 0.44

    metrics = _body_metrics(t)
    spacing = _bullet_spacing(metrics)
    y = text_top
    bottom_limit = 6.44 if t.get("footer_enabled", True) else 6.82
    for i, pt in enumerate(data.get("points", [])[:4]):
        bullet_height = _add_bullet(slide, 0.72, y, pt, t, i, width=11.7, size=max(metrics["size"] - 1, 15))
        y += bullet_height + spacing
        if y >= bottom_limit:
            break

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _summary_slide(prs, data, t, slide_index, media_assets=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, t["light"])
    _rect(slide, 0, 0, 13.33, 1.38, t["header"])
    metrics = _body_metrics(t)
    _tb(
        slide,
        0.6,
        0.2,
        12,
        1.0,
        data.get("title", "핵심 요약"),
        t.get("title_font_size", DEFAULT_TITLE_FONT_SIZE),
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        font_name=t["font_name"],
    )

    y = 1.6
    spacing = max(0.16, min(metrics["gap"] - 0.42, 0.34))
    bottom_limit = 6.45 if t.get("footer_enabled", True) else 6.8
    for i, pt in enumerate(data.get("points", [])[:6]):
        if y >= bottom_limit:
            break
        _rect(slide, 0.5, y + 0.06, 0.28, 0.28, t["confirm"])
        _tb(
            slide,
            0.5,
            y + 0.04,
            0.28,
            0.3,
            "✓",
            11,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            font_name=t["font_name"],
        )
        _tb(
            slide,
            0.94,
            y,
            11.8,
            _estimate_text_height(pt, 11.8, t.get("body_font_size", DEFAULT_BODY_FONT_SIZE), minimum=0.52),
            pt,
            t.get("body_font_size", DEFAULT_BODY_FONT_SIZE),
            color=t["text"],
            font_name=t["font_name"],
        )
        y += _estimate_text_height(pt, 11.8, t.get("body_font_size", DEFAULT_BODY_FONT_SIZE), minimum=0.52) + spacing

    _add_footer(slide, slide_index, t)
    _add_logo(slide, t.get("logo_path"))
    _set_notes(slide, data)


def _build_presentation(slides_data: list, design: dict, media_assets: list[dict] | None = None) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    t = _resolve_theme(design)
    content_handlers = {
        "classic": _content_slide_classic,
        "split": _content_slide_split,
        "card": _content_slide_card,
        "highlight": _content_slide_highlight,
        "process": _content_slide_process,
        "compare": _content_slide_compare,
        "image_left": _content_slide_image_left,
        "image_top": _content_slide_image_top,
        "chapter": _chapter_slide,
    }
    for idx, slide_data in enumerate(slides_data, start=1):
        slide_type = slide_data.get("type", "content")
        if slide_type == "title":
            _title_slide(prs, slide_data, t, idx, media_assets)
        elif slide_type == "agenda":
            _agenda_slide(prs, slide_data, t, idx, media_assets)
        elif slide_type == "summary":
            _summary_slide(prs, slide_data, t, idx, media_assets)
        elif str(slide_data.get("role") or "").lower() == "chapter":
            _chapter_slide(prs, slide_data, t, idx, media_assets)
        else:
            layout_id = _resolve_content_layout(slide_data, t)
            content_fn = content_handlers.get(layout_id, _content_slide_classic)
            content_fn(prs, slide_data, t, idx, media_assets)

    return prs


def generate_pptx(slides_data: list, output_path: str, design: dict, media_assets: list[dict] | None = None) -> str:
    prs = _build_presentation(slides_data, design, media_assets)
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs.save(output_path)
    return output_path


def generate_pptx_bytes(slides_data: list, design: dict, media_assets: list[dict] | None = None):
    from io import BytesIO
    prs = _build_presentation(slides_data, design, media_assets)
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
